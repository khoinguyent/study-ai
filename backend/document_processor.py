import os
import tempfile
from typing import List, Dict, Any
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
import openai
from dotenv import load_dotenv
import re
from langchain.text_splitter import RecursiveCharacterTextSplitter

load_dotenv()

# Configure OpenAI
openai.api_key = os.getenv('OPENAI_API_KEY')

class DocumentProcessor:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )
    
    def extract_text_from_file(self, file_path: str, file_type: str) -> str:
        """
        Extract text from various file types
        
        Args:
            file_path: Path to the file
            file_type: Type of file (pdf, docx, pptx, txt)
            
        Returns:
            Extracted text content
        """
        try:
            if file_type.lower() == 'pdf':
                return self._extract_from_pdf(file_path)
            elif file_type.lower() in ['doc', 'docx']:
                return self._extract_from_docx(file_path)
            elif file_type.lower() in ['ppt', 'pptx']:
                return self._extract_from_pptx(file_path)
            elif file_type.lower() == 'txt':
                return self._extract_from_txt(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
        except Exception as e:
            raise Exception(f"Failed to extract text from {file_path}: {str(e)}")
    
    def _extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file"""
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                if page_text.strip():
                    text += f"\n\n--- Page {page_num + 1} ---\n\n{page_text}"
        return text.strip()
    
    def _extract_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX file"""
        doc = DocxDocument(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text + "\n"
        return text.strip()
    
    def _extract_from_pptx(self, file_path: str) -> str:
        """Extract text from PPTX file"""
        prs = Presentation(file_path)
        text = ""
        for slide_num, slide in enumerate(prs.slides):
            slide_text = f"\n\n--- Slide {slide_num + 1} ---\n\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
            if slide_text.strip():
                text += slide_text
        return text.strip()
    
    def _extract_from_txt(self, file_path: str) -> str:
        """Extract text from TXT file"""
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read().strip()
    
    def chunk_text(self, text: str) -> List[Dict[str, Any]]:
        """
        Split text into chunks for processing
        
        Args:
            text: Input text to chunk
            
        Returns:
            List of chunk dictionaries with content and metadata
        """
        # Clean and preprocess text
        cleaned_text = self._clean_text(text)
        
        # Split text into chunks
        chunks = self.text_splitter.split_text(cleaned_text)
        
        # Create chunk objects with metadata
        chunk_objects = []
        for i, chunk in enumerate(chunks):
            chunk_obj = {
                'content': chunk.strip(),
                'chunk_index': i,
                'metadata': {
                    'chunk_size': len(chunk),
                    'word_count': len(chunk.split())
                }
            }
            chunk_objects.append(chunk_obj)
        
        return chunk_objects
    
    def _clean_text(self, text: str) -> str:
        """Clean and preprocess text"""
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove special characters but keep punctuation
        text = re.sub(r'[^\w\s\.\,\!\?\;\:\-\(\)\[\]\{\}]', '', text)
        
        # Normalize line breaks
        text = text.replace('\r\n', '\n').replace('\r', '\n')
        
        return text.strip()
    
    def generate_embeddings(self, text: str) -> List[float]:
        """
        Generate embeddings for text using OpenAI
        
        Args:
            text: Text to generate embeddings for
            
        Returns:
            List of embedding values
        """
        try:
            response = openai.Embedding.create(
                model="text-embedding-ada-002",
                input=text
            )
            return response['data'][0]['embedding']
        except Exception as e:
            raise Exception(f"Failed to generate embeddings: {str(e)}")
    
    def process_document(self, file_path: str, file_type: str) -> List[Dict[str, Any]]:
        """
        Complete document processing pipeline
        
        Args:
            file_path: Path to the document file
            file_type: Type of document
            
        Returns:
            List of processed chunks with embeddings
        """
        # Extract text
        text = self.extract_text_from_file(file_path, file_type)
        
        # Chunk text
        chunks = self.chunk_text(text)
        
        # Generate embeddings for each chunk
        processed_chunks = []
        for chunk in chunks:
            try:
                embedding = self.generate_embeddings(chunk['content'])
                chunk['embedding'] = embedding
                processed_chunks.append(chunk)
            except Exception as e:
                print(f"Failed to generate embedding for chunk {chunk['chunk_index']}: {e}")
                # Continue with other chunks
                continue
        
        return processed_chunks

# Global document processor instance
document_processor = DocumentProcessor() 